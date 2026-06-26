from pathlib import Path

from pptx import Presentation
from pydantic import BaseModel

from openclaw.workspace import WorkspaceError, WorkspaceService


class SlideSummary(BaseModel):
    index: int
    title: str
    shape_count: int


class PresentationSummary(BaseModel):
    path: str
    slide_count: int
    slides: list[SlideSummary]
    compatibility_note: str = ""


class PresentationResult(BaseModel):
    path: str
    saved: bool
    message: str = ""


class PresentationService:
    def __init__(self, workspace: WorkspaceService) -> None:
        self.workspace = workspace

    def _resolve_presentation(self, relative_path: str, must_exist: bool = True) -> Path:
        path = self.workspace.resolve_inside(relative_path)
        suffix = path.suffix.lower()
        if suffix == ".show":
            raise WorkspaceError(".show native editing is not available in this slice; use PPTX as the editable format")
        if suffix != ".pptx":
            raise WorkspaceError("only .pptx files are editable in this slice")
        if must_exist and not path.exists():
            raise WorkspaceError(f"presentation does not exist: {relative_path}")
        return path

    def create_presentation(self, relative_path: str, title: str, subtitle: str = "") -> PresentationResult:
        path = self._resolve_presentation(relative_path, must_exist=False)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(path)
        return PresentationResult(path=relative_path, saved=True, message="presentation created")

    def summarize_presentation(self, relative_path: str) -> PresentationSummary:
        path = self._resolve_presentation(relative_path)
        presentation = Presentation(path)
        slides: list[SlideSummary] = []
        for index, slide in enumerate(presentation.slides, start=1):
            title = ""
            if slide.shapes.title is not None:
                title = slide.shapes.title.text
            slides.append(SlideSummary(index=index, title=title, shape_count=len(slide.shapes)))
        return PresentationSummary(path=relative_path, slide_count=len(slides), slides=slides)

    def add_title_slide(self, relative_path: str, title: str, subtitle: str = "") -> PresentationResult:
        path = self._resolve_presentation(relative_path)
        presentation = Presentation(path)
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        presentation.save(path)
        return PresentationResult(path=relative_path, saved=True, message="title slide added")

    def add_bullet_slide(self, relative_path: str, title: str, bullets: list[str]) -> PresentationResult:
        path = self._resolve_presentation(relative_path)
        presentation = Presentation(path)
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
        presentation.save(path)
        return PresentationResult(path=relative_path, saved=True, message="bullet slide added")

    def show_compatibility_note(self, relative_path: str) -> PresentationSummary:
        path = self.workspace.resolve_inside(relative_path)
        if path.suffix.lower() != ".show":
            raise WorkspaceError("compatibility note is only for .show files")
        return PresentationSummary(
            path=relative_path,
            slide_count=0,
            slides=[],
            compatibility_note=".show files require Hancom Show conversion or native automation; PPTX is the current editable format.",
        )
